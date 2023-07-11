input_req = {'Upload Template':'file_upload', 'Upload Images':'file_upload'}

import argparse
import sys
from pptx import Presentation
from pptx.util import Inches

# Create an argument parser
# parser = argparse.ArgumentParser(description='Create presentation.')

# Add command-line arguments
# parser.add_argument('arg1', nargs='+')

# Parse the command-line arguments
# args = parser.parse_args()

temp_path = sys.argv[1:]

# print('main_temp_path: ', temp_path)
template = temp_path[0]
img_list = temp_path[1:]
#
# print('temp_path: ', template)
# print('img_list: ', img_list)

left = Inches(1.5)
top = Inches(0.75)
width = Inches(10.5)
height = Inches(6)

# template = st.file_uploader('Upload template')
# photos = st.file_uploader('Upload photos', accept_multiple_files=True)

# def generate_ppt(temp_path):

# print('template: ', template)
def generate_ppt(temp_path):
    presentation = Presentation(template)
    for filename in img_list:
        # if filename.endswith(".jpg") or filename.endswith(".jpeg") or filename.endswith(".png"):
        # Create a new slide
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        # Add the image to the slide
        picture = slide.shapes.add_picture(filename, left, top, width, height)
    presentation.save('modified_template.pptx')

    return 'PPT Generated'



print(generate_ppt(temp_path))
